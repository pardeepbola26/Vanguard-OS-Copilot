# ============================================================
# Vanguard OS v2 - API Server (FastAPI)
# Premium Consulting Edition — with modes, frameworks, refine
# ============================================================

from typing import Dict, Any, List
import uvicorn
import json
import re

from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel
from contextlib import asynccontextmanager
import pandas as pd
import io
import pypdf
import pdfplumber
try:
    from docx import Document  # python-docx: DOCX parsing
except ImportError:
    Document = None
    print("[startup] python-docx not installed — .docx upload will fall back to text extraction.")
except Exception as e:
    # Catch cases where the wrong `docx` package is installed (Python-2 era)
    Document = None
    print(f"[startup] docx import failed ({e}) — .docx upload disabled. Fix: pip install python-docx")
from openai import OpenAI
import os
from dotenv import load_dotenv
import time
import uuid

load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# In-memory session storage for document context
document_sessions = {}

from vanguard_agents import (
    run_vanguard_pipeline,
    run_vanguard_pipeline_stream,
    summary_agent,
    slides_agent,
    market_mapping_agent,
)
from fastapi.responses import StreamingResponse

@asynccontextmanager
async def lifespan(app: FastAPI):
    # Startup — initialize SQLite for durable mission persistence (best-effort)
    try:
        from database import init_db
        init_db()
    except Exception as e:
        print(f"DB init skipped: {e}")
    yield

app = FastAPI(
    title="Vanguard OS v2",
    description="Multi-Agent Strategy Copilot — Vanguard",
    version="2.3",
    lifespan=lifespan,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# Mount static files
app.mount("/static", StaticFiles(directory="static"), name="static")

@app.get("/")
async def read_root():
    return FileResponse("vanguard.html")


# ------------------------------------------------------------
# Pydantic Models
# ------------------------------------------------------------
class VanguardInput(BaseModel):
    situation: str
    goal: str
    constraints: str = ""
    problem_type: str = "general_business"
    mode: str = "advanced"          # "lite" | "advanced"
    frameworks: List[str] = []      # Frameworks selected in UI
    selected_agents: List[str] = [] # Agents to run (empty = all)
    numbers: str = ""               # Pasted numeric/CSV block (optional)
    document_session_id: str = ""   # Session ID from document upload
    topic: str = ""                 # For Deep Dive mode
    context: str = ""               # For Deep Dive mode


class RefineInput(BaseModel):
    situation: str
    goal: str
    constraints: str = ""
    problem_type: str = "general_business"
    mode: str = "advanced"
    frameworks: List[str] = []
    numbers: str = ""
    extra_context: str              # New user notes to append to situation


class SummaryInput(BaseModel):
    kernel: str
    options: str
    red_team: str
    recon: str


class SlidesInput(BaseModel):
    kernel: str
    drivers: str
    enright: str = ""
    frameworks: str
    options: str
    red_team: str
    recon: str
    crux: str = ""
    financial: str = ""


# ------------------------------------------------------------
# Endpoints
# ------------------------------------------------------------
@app.post("/vanguard/run")
def run_pipeline(request: VanguardInput) -> Dict[str, Any]:
    try:
        outputs = run_vanguard_pipeline(request.dict())
        # save_mission(request.dict(), outputs) - Disabled
        return outputs
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/vanguard/stream")
async def stream_pipeline(request: VanguardInput):
    try:
        # We need to capture the full output to save it, but we are streaming.
        # For simplicity in this streaming implementation, we will NOT save to DB automatically 
        # inside the stream generator because it's complex to aggregate.
        # Instead, the FRONTEND will send a save request after completion? 
        # OR we wrap the generator. Let's wrap the generator.
        
        async def generator_wrapper(inputs):
            print(f"DEBUG: Starting stream for goal: {inputs.get('goal', 'Unknown')[:30]}...")
            
            # Retrieve document context if session ID provided
            document_context = None
            session_id = inputs.get('document_session_id', '')
            if session_id and session_id in document_sessions:
                session_data = document_sessions[session_id]
                document_context = session_data['text']
                print(f"✓ Using document context from session: {session_id[:8]}... ({len(document_context)} chars)")
            else:
                if session_id:
                    print(f"⚠ Session ID {session_id[:8]}... not found or expired")
            
            # Add document_context to inputs
            inputs['document_context'] = document_context
            
            # Detect refine mode
            if 'refine_context' in inputs:
                try:
                    refine = inputs['refine_context']
                    print(f"✓ Refine mode detected with user commentary: {refine.get('user_commentary', '')[:50]}...")
                    
                    # Ensure situation exists (default to empty string if None)
                    if inputs.get('situation') is None:
                        inputs['situation'] = ''
                    
                    # Augment situation with refinement context
                    inputs['situation'] += f"\n\n[USER REFINEMENT]:\n{refine.get('user_commentary', '')}"
                    
                    # Update numbers if provided
                    if refine.get('updated_numbers'):
                        inputs['numbers'] = refine.get('updated_numbers')
                    
                    print(f"✓ Refine context applied successfully")
                except Exception as refine_error:
                    print(f"ERROR in refine mode: {refine_error}")
                    # Continue without refine context rather than failing
            
            # Immediate ping to confirm connection and unblock UI
            yield json.dumps({"type": "status", "data": "Initializing Mission..."}) + "\n"

            aggregated_outputs = {}
            full_text = ""
            
            # Use async iterator
            iterator = run_vanguard_pipeline_stream(inputs)
            async for chunk in iterator:
                yield chunk
                # Aggregate for saving
                try:
                    data = json.loads(chunk.strip())
                    if data["type"] != "done":
                        aggregated_outputs[data["type"]] = data["data"]
                except:
                    pass
            
            # Save after stream ends
            # save_mission(inputs, aggregated_outputs) - Disabled

        return StreamingResponse(
            generator_wrapper(request.dict()),
            media_type="application/x-ndjson"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/vanguard/history")
def get_mission_history():
    try:
        from database import get_history
        return get_history()
    except Exception as e:
        print(f"History load error: {e}")
        return []

@app.get("/vanguard/history/{mission_id}")
def get_mission_details(mission_id: int):
    try:
        from database import load_mission
        mission = load_mission(mission_id)
        if mission is None:
            raise HTTPException(status_code=404, detail="Mission not found")
        return mission
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


class SaveMissionInput(BaseModel):
    inputs: Dict[str, Any]
    outputs: Dict[str, Any]


@app.post("/vanguard/history")
def save_mission_endpoint(req: SaveMissionInput):
    """Persist a completed mission for durable history. Frontend posts after a run completes."""
    try:
        from database import save_mission
        mission_id = save_mission(req.inputs, req.outputs)
        return {"id": mission_id, "ok": mission_id is not None}
    except Exception as e:
        print(f"Save error: {e}")
        return {"id": None, "ok": False, "error": str(e)}


@app.post("/vanguard/upload")
async def upload_data_file(file: UploadFile = File(...)):
    """
    Upload CSV or Excel file and auto-extract financial metrics.
    Returns structured JSON with detected metrics and suggestions.
    """
    try:
        # Read file content
        contents = await file.read()
        
        # Determine file type and parse
        if file.filename.endswith('.csv'):
            df = pd.read_csv(io.BytesIO(contents))
        elif file.filename.endswith(('.xlsx', '.xls')):
            df = pd.read_excel(io.BytesIO(contents))
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type. Use CSV or Excel.")
        
        # Auto-detect metrics
        detected_metrics = {}
        suggestions = []
        
        # Common column name patterns
        revenue_patterns = ['revenue', 'sales', 'income', 'total revenue', 'gross revenue']
        cost_patterns = ['cost', 'expense', 'cogs', 'operating cost', 'opex']
        customer_patterns = ['customers', 'users', 'accounts', 'subscribers']
        cac_patterns = ['cac', 'customer acquisition cost', 'acquisition cost']
        ltv_patterns = ['ltv', 'lifetime value', 'clv', 'customer lifetime value']
        churn_patterns = ['churn', 'churn rate', 'attrition']
        
        columns_lower = [col.lower() for col in df.columns]
        
        # Detect and extract metrics
        for col, col_lower in zip(df.columns, columns_lower):
            if any(pattern in col_lower for pattern in revenue_patterns):
                detected_metrics['revenue'] = {
                    'column': col,
                    'latest': float(df[col].iloc[-1]) if len(df) > 0 else 0,
                    'growth': f"{((df[col].iloc[-1] / df[col].iloc[0] - 1) * 100):.1f}%" if len(df) > 1 else "N/A",
                    'data_points': len(df)
                }
                suggestions.append(f"Revenue: ${detected_metrics['revenue']['latest']:,.0f} (Growth: {detected_metrics['revenue']['growth']})")
            
            elif any(pattern in col_lower for pattern in cost_patterns):
                detected_metrics['costs'] = {
                    'column': col,
                    'latest': float(df[col].iloc[-1]) if len(df) > 0 else 0,
                    'data_points': len(df)
                }
                suggestions.append(f"Costs: ${detected_metrics['costs']['latest']:,.0f}")
            
            elif any(pattern in col_lower for pattern in customer_patterns):
                detected_metrics['customers'] = {
                    'column': col,
                    'latest': int(df[col].iloc[-1]) if len(df) > 0 else 0,
                    'data_points': len(df)
                }
                suggestions.append(f"Customers: {detected_metrics['customers']['latest']:,}")
            
            elif any(pattern in col_lower for pattern in cac_patterns):
                detected_metrics['cac'] = {
                    'column': col,
                    'latest': float(df[col].iloc[-1]) if len(df) > 0 else 0,
                }
                suggestions.append(f"CAC: ${detected_metrics['cac']['latest']:.2f}")
            
            elif any(pattern in col_lower for pattern in ltv_patterns):
                detected_metrics['ltv'] = {
                    'column': col,
                    'latest': float(df[col].iloc[-1]) if len(df) > 0 else 0,
                }
                suggestions.append(f"LTV: ${detected_metrics['ltv']['latest']:.2f}")
            
            elif any(pattern in col_lower for pattern in churn_patterns):
                detected_metrics['churn'] = {
                    'column': col,
                    'latest': float(df[col].iloc[-1]) if len(df) > 0 else 0,
                }
                suggestions.append(f"Churn: {detected_metrics['churn']['latest']:.1f}%")
        
        # Generate auto-fill text
        auto_fill_text = "\n".join(suggestions) if suggestions else "No standard metrics detected."
        
        return {
            "filename": file.filename,
            "rows": len(df),
            "columns": list(df.columns),
            "detected_metrics": detected_metrics,
            "auto_fill_text": auto_fill_text,
            "preview": df.head(5).to_dict('records')
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"File processing error: {str(e)}")


@app.post("/vanguard/upload-document")
async def upload_document(file: UploadFile = File(...)):
    """
    Upload PDF or DOCX document (10-K, annual report, etc.) and extract key intelligence.
    Returns structured analysis for auto-population of strategy inputs.
    """
    try:
        contents = await file.read()
        extracted_text = ""
        
        # Extract text based on file type
        if file.filename.endswith('.pdf'):
            # Try pdfplumber first (better layout preservation)
            try:
                with pdfplumber.open(io.BytesIO(contents)) as pdf:
                    text_chunks = []
                    # Limit to first 50 pages to avoid token limits
                    for page in pdf.pages[:50]:
                        text_chunks.append(page.extract_text() or "")
                    extracted_text = "\n".join(text_chunks)
            except:
                # Fallback to pypdf
                reader = pypdf.PdfReader(io.BytesIO(contents))
                text_chunks = []
                for page in reader.pages[:50]:
                    text_chunks.append(page.extract_text() or "")
                extracted_text = "\n".join(text_chunks)
        
        elif file.filename.endswith(('.docx', '.doc')):
            if Document is None:
                raise HTTPException(status_code=503, detail="DOCX parsing unavailable — install python-docx on the server (pip install python-docx).")
            doc = Document(io.BytesIO(contents))
            text_chunks = [para.text for para in doc.paragraphs]
            extracted_text = "\n".join(text_chunks)
        
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type. Use PDF or DOCX.")
        
        # Store full document text (up to 15K chars for context enrichment)
        if len(extracted_text) > 15000:
            extracted_text = extracted_text[:15000] + "\n...[Document truncated for analysis]"
        
        # GPT-4 analysis
        analysis_prompt = f"""
Analyze this business document and extract key intelligence. Return ONLY valid JSON, no markdown formatting.

Document Text (first 15k chars):
{extracted_text}

Return this exact JSON structure:
{{
  "company_name": "Company name or 'Unknown'",
  "overview": "2-3 sentence overview",
  "key_financials": {{
    "revenue": "e.g. $500M (2023)",
    "profit": "e.g. $50M net income",
    "margins": "e.g. 35% gross margin",
    "growth_rate": "e.g. 20% YoY"
  }},
  "strategic_initiatives": ["Initiative 1", "Initiative 2", "Initiative 3"],
  "risk_factors": ["Risk 1", "Risk 2", "Risk 3"],
  "competitive_position": "Brief market position description",
  "auto_fill_text": "3-4 sentence summary for Situation field"
}}

IMPORTANT: Return ONLY the JSON object. Do not wrap in markdown code blocks. If data is missing, use "Not found".
"""

        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a business analyst extracting key intelligence from company documents."},
                {"role": "user", "content": analysis_prompt}
            ],
            max_tokens=1500,
            temperature=0.2
        )
        
        analysis_raw = response.choices[0].message.content.strip()
        
        # Log raw response for debugging
        print(f"DEBUG: GPT-4 response length: {len(analysis_raw)}")
        print(f"DEBUG: First 300 chars: {analysis_raw[:300]}")
        
        # Parse JSON response (handle markdown code blocks)
        analysis = None
        try:
            # Try direct JSON parse first
            analysis = json.loads(analysis_raw)
        except:
            # Try extracting JSON from markdown code block
            import re
            json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', analysis_raw, re.DOTALL)
            if json_match:
                try:
                    analysis = json.loads(json_match.group(1))
                except:
                    pass
            
            if not analysis:
                # Try finding JSON object without markdown
                json_match = re.search(r'\{.*\}', analysis_raw, re.DOTALL)
                if json_match:
                    try:
                        analysis = json.loads(json_match.group(0))
                        print("✓ Found JSON object in text")
                    except:
                        pass
        
        # Fallback if all parsing failed
        if not analysis:
            print("✗ ALL parsing methods failed. Using fallback.")
            analysis = {
                "company_name": "Parsing Failed",
                "overview": "Could not extract structured data. Check server console for raw GPT response.",
                "key_financials": {"note": "Extraction failed"},
                "strategic_initiatives": [],
                "risk_factors": [],
                "competitive_position": "Not analyzed",
                "auto_fill_text": f"{file.filename}: Text extracted ({len(extracted_text)} chars) but GPT-4 response not parseable. Check server logs."
            }
        
        # Create session and store full document text
        session_id = str(uuid.uuid4())
        document_sessions[session_id] = {
            'text': extracted_text,
            'analysis': analysis,
            'filename': file.filename,
            'timestamp': time.time()
        }
        
        print(f"✓ Created document session: {session_id[:8]}... ({len(extracted_text)} chars)")
        
        return {
            "session_id": session_id,
            "filename": file.filename,
            "extracted_chars": len(extracted_text),
            "analysis": analysis,
            "auto_fill_text": analysis.get("auto_fill_text", ""),
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Document processing error: {str(e)}")

@app.post("/vanguard/refine")
def refine_pipeline(request: RefineInput) -> Dict[str, Any]:
    """
    Re-run the entire pipeline with additional user context appended
    so the whole plan updates (kernel, options, red team, execution, etc.).
    """
    try:
        augmented_situation = (
            request.situation
            + "\n\nAdditional context from user refinement:\n"
            + request.extra_context
        )
        inputs = {
            "situation": augmented_situation,
            "goal": request.goal,
            "constraints": request.constraints,
            "problem_type": request.problem_type,
            "mode": request.mode,
            "frameworks": request.frameworks,
            "numbers": request.numbers,
        }
        return run_vanguard_pipeline(inputs)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/vanguard/summary")
def generate_summary(request: SummaryInput):
    try:
        summary = summary_agent(
            request.kernel,
            request.options,
            request.red_team,
            request.recon,
        )
        return {"summary": summary}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/vanguard/slides")
def generate_slides(request: SlidesInput):
    try:
        slides = slides_agent(
            request.kernel,
            request.drivers,
            request.enright,
            request.frameworks,
            request.options,
            request.red_team,
            request.recon,
            request.crux,
            request.financial,
        )
        return {"slides": slides}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


if __name__ == "__main__":
    uvicorn.run("api_server:app", host="127.0.0.1", port=8000, reload=True)

# ------------------------------------------------------------
# PPTX Export Logic
# ------------------------------------------------------------
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from fastapi.responses import FileResponse
import os

@app.post("/vanguard/export/pptx")
def export_pptx(request: SlidesInput):
    if not PPTX_AVAILABLE:
        raise HTTPException(status_code=500, detail="python-pptx library not installed.")

    try:
        prs = Presentation()

        # Helper to add a title slide
        def add_title_slide(title, subtitle):
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = subtitle

        # Helper to add a content slide
        def add_content_slide(title, content):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            
            # Add text box if content is just text
            tf = slide.placeholders[1].text_frame
            tf.text = content

        # 1. Title Slide
        add_title_slide("Strategy Deck", "Generated by Vanguard OS")

        # 2. Kernel
        add_content_slide("Strategic Kernel", request.kernel)

        # 3. Drivers
        add_content_slide("Key Drivers", request.drivers)

        # 4. Frameworks
        add_content_slide("Framework Analysis", request.frameworks)

        # 5. Options
        add_content_slide("Strategic Options", request.options)

        # 6. Red Team
        add_content_slide("Red Team Analysis", request.red_team)

        # 7. Execution
        add_content_slide("Execution Plan", request.recon)

        # Save to a temporary file
        output_filename = "Vanguard_Strategy.pptx"
        output_path = os.path.abspath(output_filename)
        prs.save(output_path)
        
        # Return as a downloadable file
        return FileResponse(
            path=output_path,
            filename=output_filename,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ------------------------------------------------------------
# MARKET MAPPING
# ------------------------------------------------------------

class MarketMapRequest(BaseModel):
    industry: str
    geo_scope: str = "Global"
    segments: str = "All"

@app.post("/market_map")
async def run_market_map(request: MarketMapRequest):
    """Generate market mapping analysis"""
    result = await market_mapping_agent(
        request.industry,
        request.geo_scope,
        request.segments
    )
    return json.loads(result)

